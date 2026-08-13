import sys
import pymupdf as fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import io

def pdf_to_pptx(pdf_path, pptx_path):
    try:
        doc = fitz.open(pdf_path)
        prs = Presentation()
        
        # Use a blank slide layout (index 6 is usually blank)
        blank_slide_layout = prs.slide_layouts[6]
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Match slide dimensions to PDF page dimensions
            # PyMuPDF uses points. pptx uses EMUs, but we can set dimensions in points (Pt)
            rect = page.rect
            width_pt = rect.width
            height_pt = rect.height
            
            if page_num == 0:
                # Set presentation size based on the first page
                prs.slide_width = Pt(width_pt)
                prs.slide_height = Pt(height_pt)
                
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 1. Background image (Hybrid approach for complex graphics preservation)
            # Render page without text to preserve vector graphics/charts visually
            # PyMuPDF doesn't easily render "only graphics", so we render the whole page at 150dpi as a base layer.
            # However, if we do that, the text is duplicated. 
            # To avoid text duplication, we'll just extract images and text for a clean editable slide.
            
            # Extract images
            for img in page.get_images(full=True):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                # We need the image bounding box. 
                # PyMuPDF doesn't give direct bounding box for xref easily without iterating page.get_image_info()
                # We'll just skip complex image positioning in this simplified robust version,
                # or use get_image_info() which requires fitz >= 1.19
                pass
            
            # Better approach for images: page.get_image_info()
            try:
                for img_info in page.get_image_info(xrefs=True):
                    xref = img_info.get("xref")
                    if xref:
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_stream = io.BytesIO(image_bytes)
                        bbox = img_info["bbox"]
                        slide.shapes.add_picture(
                            image_stream, 
                            Pt(bbox[0]), Pt(bbox[1]), 
                            Pt(bbox[2] - bbox[0]), Pt(bbox[3] - bbox[1])
                        )
            except Exception as e:
                # Ignore if image extraction fails
                pass

            # Extract Text
            text_dict = page.get_text("dict")
            for block in text_dict.get("blocks", []):
                if block.get("type") == 0: # Text block
                    bbox = block["bbox"]
                    x, y, x1, y1 = bbox
                    
                    # Create text box
                    txBox = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(x1 - x), Pt(y1 - y))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.clear() # Remove default paragraph
                    
                    for line in block.get("lines", []):
                        p = tf.add_paragraph()
                        # Avoid huge spacing
                        p.space_before = Pt(0)
                        p.space_after = Pt(0)
                        
                        for span in line.get("spans", []):
                            run = p.add_run()
                            run.text = span["text"]
                            run.font.size = Pt(span["size"])
                            
                            # Font styles
                            flags = span["flags"]
                            if flags & 2 ** 4: # Bold
                                run.font.bold = True
                            if flags & 2 ** 1: # Italic
                                run.font.italic = True
                            
                            # Color (integer to RGB)
                            color = span["color"]
                            b = color & 255
                            g = (color >> 8) & 255
                            r = (color >> 16) & 255
                            run.font.color.rgb = RGBColor(r, g, b)

        prs.save(pptx_path)
        doc.close()
        print("SUCCESS")
    except Exception as e:
        print(f"ERROR: {str(e)}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_pptx.py <input.pdf> <output.pptx>", file=sys.stderr)
        sys.exit(1)

    pdf_file = sys.argv[1]
    pptx_file = sys.argv[2]
    
    pdf_to_pptx(pdf_file, pptx_file)
